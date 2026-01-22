"""
Doküman Yükleme Modülü
PDF, DOCX, PPTX ve metin dosyalarını işler
"""

from pathlib import Path
from typing import List, Optional, Tuple
import logging

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Doküman yükleyicisi - PDF, DOCX, PPTX ve web içeriğini işler"""
    
    def __init__(self):
        """Başlat"""
        self.supported_formats = ['.pdf', '.txt', '.docx', '.pptx']
    
    def load_pdf(self, file_path: str) -> Optional[str]:
        """
        PDF dosyasını yükle ve metni çıkar
        
        Args:
            file_path: PDF dosyasının yolu
            
        Returns:
            Çıkarılan metin
        """
        try:
            try:
                # Try the legacy PyPDF2 import first
                from PyPDF2 import PdfReader
            except ImportError:
                # Fall back to the modern package name
                from pypdf import PdfReader

            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                text = ""

                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    # extract_text is available on both PyPDF2 and pypdf readers
                    text += (page.extract_text() or "") + "\n"

                logger.info(f"PDF yüklendi: {file_path}")
                return text

        except ImportError:
            logger.error("PDF kütüphanesi yüklü değil. 'pip install pypdf' veya 'pip install PyPDF2' çalıştır")
            return None
        except Exception as e:
            logger.error(f"PDF yükleme hatası: {e}")
            return None
    
    def load_text(self, file_path: str) -> Optional[str]:
        """
        Metin dosyasını yükle
        
        Args:
            file_path: Metin dosyasının yolu
            
        Returns:
            Dosya içeriği
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                logger.info(f"Metin dosyası yüklendi: {file_path}")
                return text
        except Exception as e:
            logger.error(f"Metin dosyası yükleme hatası: {e}")
            return None
    
    def load_docx(self, file_path: str) -> Optional[str]:
        """
        DOCX dosyasını yükle ve metni çıkar
        
        Args:
            file_path: DOCX dosyasının yolu
            
        Returns:
            Çıkarılan metin
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            text = ""
            
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Tablolardan da metin al
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            logger.info(f"DOCX yüklendi: {file_path}")
            return text
        except ImportError:
            logger.error("python-docx kütüphanesi yüklü değil. 'pip install python-docx' çalıştır")
            return None
        except Exception as e:
            logger.error(f"DOCX yükleme hatası: {e}")
            return None
    
    def load_pptx(self, file_path: str) -> Optional[str]:
        """
        PPTX dosyasını yükle ve metni çıkar
        
        Args:
            file_path: PPTX dosyasının yolu
            
        Returns:
            Çıkarılan metin
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"--- SLIDE {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            logger.info(f"PPTX yüklendi: {file_path}")
            return text
        except ImportError:
            logger.error("python-pptx kütüphanesi yüklü değil. 'pip install python-pptx' çalıştır")
            return None
        except Exception as e:
            logger.error(f"PPTX yükleme hatası: {e}")
            return None
    
    def load_document(self, file_path: str) -> Optional[str]:
        """
        Desteklenen herhangi bir dosyayı yükle
        
        Args:
            file_path: Dosyanın yolu
            
        Returns:
            Dosya içeriği
        """
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        if suffix == '.pdf':
            return self.load_pdf(file_path)
        elif suffix == '.txt':
            return self.load_text(file_path)
        elif suffix == '.docx':
            return self.load_docx(file_path)
        elif suffix == '.pptx':
            return self.load_pptx(file_path)
        else:
            logger.error(f"Desteklenmeyen dosya formatı: {suffix}")
            return None


class TextSplitter:
    """Metni parçalara böl - RAG için hazırlık"""
    
    def __init__(self, chunk_size: int = 1000, overlap: int = 200):
        """
        Başlat
        
        Args:
            chunk_size: Her parçanın maksimum karakteri
            overlap: Parçalar arasındaki örtüşme
        """
        self.chunk_size = chunk_size
        self.overlap = overlap
    
    def split_text(self, text: str) -> List[str]:
        """
        Metni parçalara böl
        
        Args:
            text: Bölünecek metin
            
        Returns:
            Metin parçaları listesi
        """
        total_len = len(text)

        # Büyük dokümanlarda chunk/overlap değerini dinamik büyüt
        dyn_chunk = self.chunk_size
        dyn_overlap = self.overlap
        if total_len > 200_000:
            dyn_chunk = 2200
            dyn_overlap = 260
        elif total_len > 120_000:
            dyn_chunk = 1800
            dyn_overlap = 230
        elif total_len > 60_000:
            dyn_chunk = 1400
            dyn_overlap = 210
        else:
            dyn_chunk = self.chunk_size
            dyn_overlap = self.overlap

        chunks = []
        start = 0
        
        while start < total_len:
            end = start + dyn_chunk
            
            # Son kısım için word boundary'yi bul
            if end < total_len:
                end = text.rfind(' ', start, end)
                if end == -1 or end <= start:
                    end = start + dyn_chunk
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Sonraki parça overlap kadar geri gitsin
            start = end - dyn_overlap
        
        logger.info(f"Metin {len(chunks)} parçaya bölündü (chunk={dyn_chunk}, overlap={dyn_overlap}, len={total_len})")
        return chunks
